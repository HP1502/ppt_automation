import shelve
from pptx import Presentation
from pptx.util import Inches,Pt
import datetime
import itertools
import os
import pandas as pd
presentation = Presentation(r'template.pptx')
image1=r'image1.png'
image2=r'image2.jpg'
page_id={i+1:slide.slide_id for i,slide in enumerate(presentation.slides)}
slide1=presentation.slides.get(page_id[4])
pic=slide1.shapes.add_picture(image1,left=Inches(0),top=Inches(0),width=Inches(13.333),height=Inches(7.5))
slide2=presentation.slides.get(page_id[5]) 
pic=slide2.shapes.add_picture(image2,left=Inches(0),top=Inches(0),width=Inches(13.333),height=Inches(7.5))
l= [name for name in os.listdir(".") if os.path.isdir(name) and name.endswith('_094813')]

for i in range(len(l)):
    l[i]="".join(itertools.takewhile(lambda x: x!="_",l[i]))
#print(l)
format='%Y%m%d'
dates=[]
d_date=datetime.date(2023,1,1)
i_date=datetime.date(2023,1,1)
for i in range(len(l)):
    dates.append(datetime.datetime.strptime(l[i],format))

for i in range(len(dates)):
    if(dates[i]<dates[i+1]):
        i_date=dates[i]
        d_date=dates[i+1]
        break
d_date1=d_date.strftime("%d")+ " " + d_date.strftime("%B")+ " " + d_date.strftime("%Y")
i_date1=i_date.strftime("%d")+" " + i_date.strftime("%B")+" " + i_date.strftime("%Y")
d_slide=presentation.slides.get(page_id[1])
i_slide=presentation.slides.get(page_id[3])
tx1=d_slide.shapes.add_textbox(left=Inches(7.25), top=Inches(0.64), width=Inches(2.25), height=Inches(0.88))
tf1=tx1.text_frame
tf1.text="Delivery Date: "
tf1.text+=d_date1
tx2=i_slide.shapes.add_textbox(left=Inches(4.88), top=Inches(6.24), width=Inches(2.12), height=Inches(0.39))
tf2=tx2.text_frame
tf2.text=i_date1
tx3=i_slide.shapes.add_textbox(left=Inches(5.38), top=Inches(6.63), width=Inches(2.12), height=Inches(0.34))
tf3=tx3.text_frame
#tx3.text="tbd"
pd.read_csv('stats_others.csv', header=None).T.to_csv('output.csv', header=False, index=False)
df=pd.read_csv('output.csv')
#print(df['Area with cloud (hectares and %)'])
p_list=df['Area with cloud (hectares and %)']
percentage=str(p_list[1])
#print(percentage)
tf3.text=percentage
presentation.save(r'result.pptx')





